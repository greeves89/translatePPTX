#!/usr/bin/env python3
import sys
import os
import types
import warnings
import concurrent.futures
import threading
import time
from dotenv import load_dotenv
from tqdm import tqdm

# .env laden
load_dotenv()

# Optionen aus der .env
use_local_translation = os.getenv("LOCAL_TRANSLATION", "False").lower() == "true"
deepl_api_key = os.getenv("DEEPL_API_KEY")
use_deepl = bool(deepl_api_key and deepl_api_key.strip())

# Dummy-Injektion für das cgi-Modul in Python 3.13
if sys.version_info >= (3, 13):
    warnings.warn("cgi module removed in Python 3.13; injecting dummy module with parse_header", DeprecationWarning)
    def parse_header(value):
        parts = value.split(";")
        key = parts[0].strip()
        params = {}
        for param in parts[1:]:
            if "=" in param:
                k, v = param.split("=", 1)
                params[k.strip()] = v.strip()
        return key, params
    cgi_dummy = types.ModuleType("cgi")
    cgi_dummy.parse_header = parse_header
    sys.modules["cgi"] = cgi_dummy

# RateLimiter: maximal 2 Anfragen pro Sekunde
class RateLimiter:
    def __init__(self, rate):
        self.rate = rate
        self.lock = threading.Lock()
        self.last_call = 0
    def wait(self):
        with self.lock:
            now = time.time()
            elapsed = now - self.last_call
            wait_time = max(0, 1.0 / self.rate - elapsed)
            if wait_time > 0:
                time.sleep(wait_time)
            self.last_call = time.time()

rate_limiter = RateLimiter(rate=2)

# Übersetzungsdienste: Falls Deepl API-Key vorhanden ist, wird Deepl genutzt, ansonsten Google Translate.
if use_deepl:
    import deepl
    translator_deepl = deepl.Translator(deepl_api_key)
    print("[INFO] DeepL API-Key gefunden. Nutze DeepL für die Übersetzung.")
else:
    from googletrans import Translator as GoogleTranslator
    translator_google = GoogleTranslator()
    print("[INFO] Kein DeepL API-Key gefunden. Nutze Google Translate für die Übersetzung.")

# Lade lokale Übersetzung (M2M-100) unconditionally für Fallback und für den lokalen Modus.
from transformers import M2M100ForConditionalGeneration, M2M100Tokenizer
from langdetect import detect
print("[INFO] Lade M2M-100 Modell für lokale Übersetzung (Fallback) ...")
model_name = "facebook/m2m100_418M"
tokenizer_local = M2M100Tokenizer.from_pretrained(model_name)
model_local = M2M100ForConditionalGeneration.from_pretrained(model_name)

# Funktion für lokale Übersetzung via M2M-100 (mit automatischer Spracherkennung)
def translate_with_local(text: str, target_language: str) -> str:
    try:
        # Auto-Detect der Quellsprache
        from langdetect import detect
        src_lang = detect(text)
    except Exception as e:
        print(f"[WARNING] Automatische Spracherkennung fehlgeschlagen: {e}")
        src_lang = "en"  # Fallback: Englisch
    try:
        tokenizer_local.src_lang = src_lang
        encoded = tokenizer_local(text, return_tensors="pt", truncation=True)
        # Für M2M-100 muss forced_bos_token_id gesetzt werden.
        target_code = target_language.lower()  # M2M100 erwartet zumeist ISO-Codes in Kleinbuchstaben
        forced_bos_token_id = tokenizer_local.get_lang_id(target_code)
        generated_tokens = model_local.generate(**encoded, forced_bos_token_id=forced_bos_token_id)
        translated = tokenizer_local.decode(generated_tokens[0], skip_special_tokens=True)
        return translated
    except Exception as e:
        print(f"[ERROR] Lokale Übersetzung (M2M-100) fehlgeschlagen für '{text}': {e}")
        return text

# Funktion, die Übersetzungen mit Wiederholungslogik durchführt
def translate_text_retry(text: str, target_language: str, max_retries=3, initial_delay=1) -> str:
    delay = initial_delay
    for attempt in range(1, max_retries + 1):
        rate_limiter.wait()
        try:
            # Falls lokaler Modus aktiviert ist, nutze lokal
            if use_local_translation:
                return translate_with_local(text, target_language)
            else:
                # Online Übersetzung: Wenn Deepl vorhanden, nutze Deepl, ansonsten Google
                deepl_language_map = {
                    "en": "EN-GB",
                    "de": "DE",
                    "fr": "FR",
                    "es": "ES",
                    "it": "IT",
                    "nl": "NL",
                    "pt": "PT-PT",
                    "sv": "SV",
                    "da": "DA",
                    "fi": "FI",
                    "no": "NO",
                    "pl": "PL"
                }
                if use_deepl:
                    target_code = deepl_language_map.get(target_language.lower(), target_language.upper())
                    result = translator_deepl.translate_text(text, target_lang=target_code)
                    return result.text
                else:
                    result = translator_google.translate(text, dest=target_language)
                    return result.text
        except Exception as e:
            err_str = str(e)
            if "Too many requests" in err_str and attempt < max_retries:
                print(f"[WARNING] Übersetzungsversuch {attempt} für '{text[:30]}...' fehlgeschlagen (Too many requests). Warte {delay} Sekunden und versuche es erneut.")
                time.sleep(delay)
                delay *= 2
            else:
                print(f"[ERROR] Online-Übersetzung fehlgeschlagen für '{text}': {e}")
    # Fallback: Immer lokale Übersetzung versuchen
    print(f"[WARNING] Fallback auf lokale Übersetzung (M2M-100) für '{text[:30]}...'")
    return translate_with_local(text, target_language)

def translate_text(text: str, target_language: str) -> str:
    return translate_text_retry(text, target_language)

# Hilfsfunktion für parallele Übersetzung eines Textabschnitts
def translate_task(text: str, target_language: str):
    return translate_text(text, target_language)

# PPTX-Übersetzung (parallelisiert mit Fortschrittsbalken und Fehlersammlung)
def translate_pptx(input_path: str, target_language: str, output_path: str = None):
    from pptx import Presentation
    print(f"[INFO] Starte Übersetzung der PPTX-Datei: {input_path}")
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Die Datei wurde nicht gefunden: {input_path}")
    prs = Presentation(input_path)
    tasks = []
    error_list = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.has_text_frame:
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    for run_index, run in enumerate(paragraph.runs, start=1):
                        orig = run.text
                        if orig.strip():
                            tasks.append({
                                "run": run,
                                "text": orig,
                                "slide": slide_index,
                                "shape": shape_index,
                                "paragraph": paragraph_index,
                                "run_idx": run_index
                            })
    print(f"[INFO] Übersetze {len(tasks)} Textabschnitte in PPTX...")
    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        future_to_info = {executor.submit(translate_task, task["text"], target_language): task for task in tasks}
        for future in tqdm(concurrent.futures.as_completed(future_to_info),
                           total=len(future_to_info),
                           desc="PPTX: Übersetze Textabschnitte"):
            info = future_to_info[future]
            try:
                translated = future.result()
                info["run"].text = translated
            except Exception as e:
                err_msg = f"Folie {info['slide']}, Shape {info['shape']}, Absatz {info['paragraph']}, Run {info['run_idx']}: {info['text']} -> {e}"
                print(f"[ERROR] Übersetzung fehlgeschlagen: {err_msg}")
                error_list.append(err_msg)
    if not output_path:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_translated{ext}"
    prs.save(output_path)
    print(f"[INFO] Übersetzte PPTX-Datei gespeichert unter: {output_path}")
    if error_list:
        print("\n[SUMMARY] Fehler bei der PPTX-Übersetzung:")
        for err in error_list:
            print(f"  {err}")
    else:
        print("[SUMMARY] Alle Textabschnitte in PPTX wurden erfolgreich übersetzt.")

# DOCX-Übersetzung (parallelisiert mit Fortschrittsbalken und Fehlersammlung)
def translate_docx(input_path: str, target_language: str, output_path: str = None):
    from docx import Document
    print(f"[INFO] Starte Übersetzung der DOCX-Datei: {input_path}")
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Die Datei wurde nicht gefunden: {input_path}")
    doc = Document(input_path)
    tasks = [para for para in doc.paragraphs if para.text.strip()]
    print(f"[INFO] Übersetze {len(tasks)} Absätze in DOCX...")
    error_list = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        future_to_para = {executor.submit(translate_task, para.text, target_language): para for para in tasks}
        for future in tqdm(concurrent.futures.as_completed(future_to_para),
                           total=len(future_to_para),
                           desc="DOCX: Übersetze Absätze"):
            para = future_to_para[future]
            try:
                para.text = future.result()
            except Exception as e:
                err_msg = f"Absatz: {para.text} -> {e}"
                print(f"[ERROR] Übersetzung fehlgeschlagen: {err_msg}")
                error_list.append(err_msg)
    if not output_path:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_translated{ext}"
    doc.save(output_path)
    print(f"[INFO] Übersetzte DOCX-Datei gespeichert unter: {output_path}")
    if error_list:
        print("\n[SUMMARY] Fehler bei der DOCX-Übersetzung:")
        for err in error_list:
            print(f"  {err}")
    else:
        print("[SUMMARY] Alle Absätze in DOCX wurden erfolgreich übersetzt.")

# XLSX-Übersetzung (parallelisiert mit Fortschrittsbalken und Fehlersammlung)
def translate_xlsx(input_path: str, target_language: str, output_path: str = None):
    from openpyxl import load_workbook
    print(f"[INFO] Starte Übersetzung der XLSX-Datei: {input_path}")
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Die Datei wurde nicht gefunden: {input_path}")
    wb = load_workbook(input_path)
    error_list = []
    for ws in wb.worksheets:
        print(f"[INFO] Verarbeite Arbeitsblatt: {ws.title}")
        tasks = [cell for row in ws.iter_rows() for cell in row if cell.value and isinstance(cell.value, str)]
        print(f"[INFO] Übersetze {len(tasks)} Zellen in '{ws.title}'...")
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            future_to_cell = {executor.submit(translate_task, cell.value, target_language): cell for cell in tasks}
            for future in tqdm(concurrent.futures.as_completed(future_to_cell),
                               total=len(future_to_cell),
                               desc=f"XLSX: Übersetze Zellen in {ws.title}"):
                cell = future_to_cell[future]
                try:
                    cell.value = future.result()
                except Exception as e:
                    err_msg = f"Zelle '{cell.value}' -> {e}"
                    print(f"[ERROR] Übersetzung fehlgeschlagen: {err_msg}")
                    error_list.append(err_msg)
    if not output_path:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_translated{ext}"
    wb.save(output_path)
    print(f"[INFO] Übersetzte XLSX-Datei gespeichert unter: {output_path}")
    if error_list:
        print("\n[SUMMARY] Fehler bei der XLSX-Übersetzung:")
        for err in error_list:
            print(f"  {err}")
    else:
        print("[SUMMARY] Alle Zellen in XLSX wurden erfolgreich übersetzt.")

# PDF-Übersetzung (Text extrahieren, übersetzen, Ausgabe als TXT)
def translate_pdf(input_path: str, target_language: str, output_path: str = None):
    from pdfminer.high_level import extract_text
    print(f"[INFO] Starte Übersetzung der PDF-Datei: {input_path}")
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Die Datei wurde nicht gefunden: {input_path}")
    try:
        text = extract_text(input_path)
    except Exception as e:
        print(f"[ERROR] Fehler beim Extrahieren des PDF-Texts: {e}")
        return
    print("[INFO] PDF-Text extrahiert.")
    translated_text = translate_text(text, target_language)
    if not output_path:
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_translated.txt"
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(translated_text)
        print(f"[INFO] Übersetzte PDF-Ausgabe als TXT gespeichert unter: {output_path}")
    except Exception as e:
        print(f"[ERROR] Fehler beim Speichern der übersetzten PDF-Ausgabe: {e}")

# DOC-Übersetzung (alte Word-Dokumente; Text extrahieren via textract, Ausgabe als TXT)
def translate_doc(input_path: str, target_language: str, output_path: str = None):
    import textract
    print(f"[INFO] Starte Übersetzung der DOC-Datei: {input_path}")
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Die Datei wurde nicht gefunden: {input_path}")
    try:
        text = textract.process(input_path).decode("utf-8")
    except Exception as e:
        print(f"[ERROR] Fehler beim Extrahieren des DOC-Texts: {e}")
        return
    print("[INFO] DOC-Text extrahiert.")
    translated_text = translate_text(text, target_language)
    if not output_path:
        base, _ = os.path.splitext(input_path)
        output_path = f"{base}_translated.txt"
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(translated_text)
        print(f"[INFO] Übersetzte DOC-Ausgabe als TXT gespeichert unter: {output_path}")
    except Exception as e:
        print(f"[ERROR] Fehler beim Speichern der übersetzten DOC-Ausgabe: {e}")

def translate_file(input_path: str, target_language: str, output_path: str = None):
    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".pptx":
        translate_pptx(input_path, target_language, output_path)
    elif ext == ".docx":
        translate_docx(input_path, target_language, output_path)
    elif ext == ".xlsx":
        translate_xlsx(input_path, target_language, output_path)
    elif ext == ".pdf":
        translate_pdf(input_path, target_language, output_path)
    elif ext == ".doc":
        translate_doc(input_path, target_language, output_path)
    else:
        print(f"[ERROR] Dateityp {ext} wird nicht unterstützt.")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python translate_file.py <input_file> <target_language>")
        sys.exit(1)
    input_file = sys.argv[1]
    target_lang = sys.argv[2]
    print(f"[INFO] Aufruf mit Datei: {input_file} und Zielsprache: {target_lang}")
    translate_file(input_file, target_lang)